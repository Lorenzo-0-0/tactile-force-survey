#!/usr/bin/env python3
"""Extract the Fig.2 framework structure (boxes + reference lists) from the
survey's PPTX source and emit data/framework.json (frozen schema).

Deterministic and re-runnable. Also dumps tools/framework.draft.json with the
raw extracted (label, refs, x, y) tuples for audit.

Fails loudly (exit 1) when:
  * a canonical box cannot be matched to any PPTX label, or
  * a ref-bearing text box cannot be assigned to a canonical box.
"""

import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX_PATH = (
    "/Users/lijingliang/Library/CloudStorage/OneDrive-NanyangTechnologicalUniversity/"
    "Tactile:Force-aware Robot Learning综述/draw main figure ver2.pptx"
)
SITE_ROOT = Path(__file__).resolve().parent.parent
OUT_JSON = SITE_ROOT / "data" / "framework.json"
DRAFT_JSON = SITE_ROOT / "tools" / "framework.draft.json"

REF_TOKEN = re.compile(r"\[\d+\]")
REF_NUM = re.compile(r"\[(\d+)\]")
HEADER_RE = re.compile(r"^\(([a-h])\)\s+(.*)$", re.S)

BAND_TITLES = {
    "fusion-phase forwarded",
    "phase 1 predicted modalities",
    "phase 2 predicted modalities",
    "phase 1 forwarded",
}

EMU = 1  # all coordinates kept in EMU


# --------------------------------------------------------------------------
# 1. Flatten the shape tree to absolute (slide-space) coordinates.
# --------------------------------------------------------------------------

def flatten_shapes(shapes, transform, group_path):
    """Yield (shape, abs_bbox, group_path) for every non-group shape.

    `transform` maps child-space (x, y) -> absolute and carries scale factors:
    (fx, fy, ox, oy) so that abs = (ox + x * fx, oy + y * fy).
    """
    fx, fy, ox, oy = transform
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            x = sh.element.grpSpPr.xfrm
            if x is None or x.chExt is None or x.chExt.cx == 0 or x.chExt.cy == 0:
                # degenerate group: treat children as untransformed
                sub = transform
            else:
                # group's own off/ext are in *parent* space -> absolutize first
                g_ox = ox + x.off.x * fx
                g_oy = oy + x.off.y * fy
                g_fx = fx * x.ext.cx / x.chExt.cx
                g_fy = fy * x.ext.cy / x.chExt.cy
                sub = (g_fx, g_fy, g_ox - x.chOff.x * g_fx, g_oy - x.chOff.y * g_fy)
            yield from flatten_shapes(sh.shapes, sub, group_path + (sh.shape_id,))
        else:
            if sh.left is None:
                continue
            ax = ox + sh.left * fx
            ay = oy + sh.top * fy
            aw = (sh.width or 0) * fx
            ah = (sh.height or 0) * fy
            yield sh, (ax, ay, aw, ah), group_path


def is_pure_reflist(text):
    if not REF_TOKEN.search(text):
        return False
    leftover = REF_TOKEN.sub("", text)
    return re.sub(r"[\s,;.]+", "", leftover) == ""


def norm(text):
    return re.sub(r"\s+", " ", text).strip().lower()


# --------------------------------------------------------------------------
# 2. Canonical inventory (FROZEN ids/panels; refs filled from PPTX).
# --------------------------------------------------------------------------

MOD_LABEL = {
    "V": "Vision", "L": "Language", "F": "Force–Torque", "T": "Tactile",
    "S": "State", "A": "Action Chunk", "Fc": "Reference Force",
    "K": "Control Stiffness", "Other": "Others",
}

COMPONENTS = [
    {"id": "perception", "letter": "a", "title": "Multi-modality Perception",
     "kind": "main", "panel": "fusion",
     "boxes": [("v", "V", "Vision"), ("l", "L", "Language"),
               ("f", "F", "Force–Torque"), ("t", "T", "Tactile"),
               ("s", "S", "State"), ("others", "others", "Others")]},
    {"id": "fusion", "letter": "b", "title": "Multi-modality Fusion & Encoding",
     "kind": "main", "panel": "fusion",
     "boxes": [("vf", "VF", "VF"), ("vt", "VT", "VT"), ("lt", "LT", "LT"),
               ("fs", "FS", "FS"), ("vlt", "VLT", "VLT"), ("vfs", "VFS", "VFS"),
               ("vts", "VTS", "VTS"), ("vlfs", "VLFS", "VLFS"),
               ("vlts", "VLTS", "VLTS"), ("vfts", "VFTS", "VFTS"),
               ("vlfts", "VLFTS", "VLFTS"), ("others", "others", "Others")]},
    {"id": "reconstruction", "letter": "c", "title": "Reconstruction",
     "kind": "side", "panel": "phase1",
     "boxes": [("t", "T", "Tactile"), ("v", "V", "Vision"),
               ("f", "F", "Force–Torque"), ("l", "L", "Language"),
               ("s", "S", "State")]},
    {"id": "policy1", "letter": "d", "title": "Primary Policy",
     "kind": "main", "panel": "phase1", "phaseLabel": "Phase 1",
     "boxes": [("vla", "vla", "Vision-Language-Action (VLA)"),
               ("dp", "dp", "Diffusion · Flow Matching"),
               ("act", "act", "Action Chunking Transformer (ACT)"),
               ("rl", "rl", "Reinforcement Learning (RL)"),
               ("reg", "reg", "Regression (MLP)")]},
    {"id": "intermediate", "letter": "f", "title": "Intermediate Modalities",
     "kind": "main", "panel": "phase1",
     "boxes": [("a", "A", "Action Chunk (A)"), ("k", "K", "Control Stiffness (K)"),
               ("fc", "Fc", "Reference Force (Fc)"), ("others", "others", "Others")]},
    {"id": "obsprediction", "letter": "e", "title": "Observation Prediction",
     "kind": "side", "panel": "phase2",
     "boxes": [("f", "F", "Force–Torque"), ("s", "S", "State"),
               ("t", "T", "Tactile")]},
    {"id": "policy2", "letter": "g", "title": "Refinement Policy",
     "kind": "main", "panel": "phase2", "phaseLabel": "Phase 2",
     "boxes": [("dp", "dp", "Diffusion · Flow Matching"), ("act", "act", "ACT"),
               ("rl", "rl", "Reinforcement Learning (RL)"),
               ("reg", "reg", "Regression (MLP)"),
               ("model", "model", "Model-based Algorithms")]},
    {"id": "control", "letter": "h", "title": "Robot-end Control",
     "kind": "main", "panel": "phase3", "phaseLabel": "Phase 3",
     "boxes": [("impedance", "impedance", "Impedance Control"),
               ("admittance", "admittance", "Admittance Control"),
               ("hybrid", "hybrid", "Hybrid Force-Position Control"),
               ("pid", "pid", "PID (PD) Control"),
               ("robotdep", "robotdep", "Robot-Dependent"),
               ("others", "others", "Others (FDCC · Bilateral)")]},
]

BANDS = [
    {"id": "band.p2in", "panel": "phase2",
     "groups": [
         {"id": "band.p2in.fwd", "title": "Fusion-phase Forwarded",
          "boxes": [("v", "V"), ("f", "F"), ("t", "T"), ("s", "S"),
                    ("other", "Other")]},
         {"id": "band.p2in.pred", "title": "Phase 1 Predicted Modalities",
          "boxes": [("a", "A"), ("fc", "Fc"), ("other", "Other")]},
     ]},
    {"id": "band.p3in", "panel": "phase3",
     "groups": [
         {"id": "band.p3in.fwd", "title": "Fusion-phase Forwarded",
          "boxes": [("f", "F"), ("s", "S")]},
         {"id": "band.p3in.pred", "title": "Phase 2 Predicted Modalities",
          "boxes": [("a", "A"), ("fc", "Fc"), ("other", "Other")]},
         {"id": "band.p3in.p1fwd", "title": "Phase 1 Forwarded",
          "boxes": [("a", "A"), ("fc", "Fc"), ("k", "K")]},
     ]},
]


# --------------------------------------------------------------------------
# 3. Bucket + per-bucket label -> canonical box id resolution.
# --------------------------------------------------------------------------

FUSION_CODES = {"vf", "vt", "lt", "fs", "vlt", "vfs", "vts", "vlfs",
                "vlts", "vfts", "vlfts"}
X_RIGHT = 11_900_000     # right side panel ((e) observation prediction)
X_LEFT_MAX = 4_200_000   # left side panel ((c) reconstruction)
X_BAND_SPLIT1 = 6_200_000
X_BAND_SPLIT2 = 9_600_000


def make_bucketer(geo):
    yb, yd = geo["b"], geo["d"]
    ye, yf, yg, yh = geo["e"], geo["f"], geo["g"], geo["h"]
    yc = geo["c"]
    y_p2band, y_p3band = geo["p2band"], geo["p3band"]
    y_d_min = min(yc, yd)

    def bucket(cx, cy):
        if cy < yb:
            return "perception"
        if cy < y_d_min:
            return "fusion"
        if cx > X_RIGHT and ye <= cy < yg:
            return "obsprediction"
        if cx < X_LEFT_MAX and cy < y_p2band:
            # left side panel ((c) Reconstruction) extends below the (f) header
            return "reconstruction"
        if cy < yf:
            return "policy1"
        if cy < y_p2band:
            return "intermediate"
        if cy < yg:
            return "band.p2in"
        if cy < y_p3band:
            return "policy2"
        if cy < yh:
            return "band.p3in"
        return "control"

    return bucket


def resolve_box(bucket_name, text, cx):
    """Map a (bucket, label-text, center-x) to a canonical box id, or None."""
    t = norm(text)
    if bucket_name == "perception":
        m = {"vision (v)": "v", "language (l)": "l", "force/torque (f)": "f",
             "tactile (t)": "t", "state (s)": "s", "(others)": "others"}
        key = m.get(t)
        return f"perception.{key}" if key else None
    if bucket_name == "fusion":
        if t in FUSION_CODES:
            return f"fusion.{t}"
        if t == "others":
            return "fusion.others"
        return None
    if bucket_name == "reconstruction":
        return f"reconstruction.{t}" if t in {"t", "v", "f", "l", "s"} else None
    if bucket_name == "policy1":
        for pat, key in [("vision-language", "vla"), ("diffusion", "dp"),
                         ("chunking", "act"), ("reinforcement", "rl"),
                         ("regression", "reg")]:
            if pat in t:
                return f"policy1.{key}"
        return None
    if bucket_name == "intermediate":
        for pat, key in [("action chunk", "a"), ("stiffness", "k"),
                         ("reference force", "fc"), ("others", "others")]:
            if pat in t:
                return f"intermediate.{key}"
        return None
    if bucket_name == "obsprediction":
        return f"obsprediction.{t}" if t in {"f", "s", "t"} else None
    if bucket_name == "band.p2in":
        grp = "fwd" if cx < X_BAND_SPLIT1 else "pred"
        valid = {"fwd": {"v", "f", "t", "s", "other"},
                 "pred": {"a", "fc", "other"}}[grp]
        return f"band.p2in.{grp}.{t}" if t in valid else None
    if bucket_name == "policy2":
        for pat, key in [("diffusion", "dp"), ("chunking", "act"),
                         ("reinforcement", "rl"), ("regression", "reg"),
                         ("model-based", "model")]:
            if pat in t:
                return f"policy2.{key}"
        return None
    if bucket_name == "band.p3in":
        grp = ("fwd" if cx < X_BAND_SPLIT1
               else "pred" if cx < X_BAND_SPLIT2 else "p1fwd")
        valid = {"fwd": {"f", "s"}, "pred": {"a", "fc", "other"},
                 "p1fwd": {"a", "fc", "k"}}[grp]
        return f"band.p3in.{grp}.{t}" if t in valid else None
    if bucket_name == "control":
        for pat, key in [("impedance", "impedance"), ("admittance", "admittance"),
                         ("hybrid", "hybrid"), ("pid", "pid"),
                         ("robot-dependent", "robotdep")]:
            if pat in t:
                return f"control.{key}"
        if t in {"others", "fdcc", "bilateral"}:
            return "control.others"
        return None
    return None


# --------------------------------------------------------------------------
# 4. Main extraction.
# --------------------------------------------------------------------------

def main():
    prs = Presentation(PPTX_PATH)
    slide = prs.slides[0]

    labels = []    # dicts: text, x, y, w, h, cx, cy, group(innermost id or None)
    refboxes = []  # dicts: text, refs, x, y, w, h, cx, cy, group
    headers = {}   # letter -> dict(text, y)
    band_title_ys = []

    for sh, (ax, ay, aw, ah), gpath in flatten_shapes(
            slide.shapes, (1.0, 1.0, 0.0, 0.0), ()):
        if not sh.has_text_frame:
            continue
        text = sh.text_frame.text.strip()
        if not text:
            continue
        rec = {
            "text": re.sub(r"\s+", " ", text),
            "x": round(ax), "y": round(ay), "w": round(aw), "h": round(ah),
            "cx": round(ax + aw / 2), "cy": round(ay + ah / 2),
            "group": gpath[-1] if gpath else None,
        }
        if is_pure_reflist(text):
            rec["refs"] = [int(n) for n in REF_NUM.findall(text)]
            refboxes.append(rec)
            continue
        hm = HEADER_RE.match(norm(text))
        if hm and hm.group(1) in "abcdefgh":
            headers[hm.group(1)] = {"text": rec["text"], "y": rec["y"]}
            continue
        if norm(text) in BAND_TITLES:
            band_title_ys.append((norm(text), rec["y"]))
            continue
        labels.append(rec)

    missing_headers = [c for c in "abcdefgh" if c not in headers]
    if missing_headers:
        sys.exit(f"FATAL: missing (x) header boxes in PPTX: {missing_headers}")

    fwd_ys = sorted(y for t, y in band_title_ys if t == "fusion-phase forwarded")
    if len(fwd_ys) != 2:
        sys.exit(f"FATAL: expected 2 'Fusion-phase Forwarded' band titles, "
                 f"found {len(fwd_ys)}")
    geo = {c: headers[c]["y"] for c in "abcdefgh"}
    geo["p2band"], geo["p3band"] = fwd_ys
    bucket = make_bucketer(geo)

    # ---- assign each ref box to a label -------------------------------
    def x_overlap(a, b):
        return a["x"] < b["x"] + b["w"] and b["x"] < a["x"] + a["w"]

    def nearest_label_above(rb, candidates):
        above = [lb for lb in candidates
                 if x_overlap(rb, lb) and lb["cy"] <= rb["cy"]]
        if not above:
            return None
        return max(above, key=lambda lb: lb["cy"])  # closest above

    problems = []
    for rb in refboxes:
        target = None
        if rb["group"] is not None:
            same = [lb for lb in labels if lb["group"] == rb["group"]]
            if len(same) == 1:
                target = same[0]
            elif len(same) > 1:
                target = nearest_label_above(rb, same)
        if target is None:
            target = nearest_label_above(rb, labels)
        if target is None:
            problems.append(f"ref box {rb['text'][:60]!r} at "
                            f"({rb['x']},{rb['y']}) has no label")
            continue
        rb["label"] = target["text"]
        rb["label_xy"] = (target["x"], target["y"])
        target.setdefault("_refs", []).extend(rb["refs"])
        rb["_target"] = target

    # ---- bucket labels and resolve to canonical boxes ------------------
    box_refs = {}      # canonical box id -> list of refs
    box_sources = {}   # canonical box id -> list of source label texts
    ignored = []
    for lb in labels:
        b = bucket(lb["cx"], lb["cy"])
        lb["bucket"] = b
        canon = resolve_box(b, lb["text"], lb["cx"])
        lb["canonical"] = canon
        if canon is None:
            ignored.append(lb)
            if lb.get("_refs"):
                problems.append(
                    f"label {lb['text'][:60]!r} (bucket {b}) carries refs "
                    f"{sorted(lb['_refs'])} but matches no canonical box")
            continue
        box_refs.setdefault(canon, []).extend(lb.get("_refs", []))
        box_sources.setdefault(canon, []).append(lb["text"])

    # ---- build output on the canonical inventory -----------------------
    all_box_ids = []
    components_out = []
    for comp in COMPONENTS:
        boxes = []
        for suffix, abbrev, label in comp["boxes"]:
            bid = f"{comp['id']}.{suffix}"
            all_box_ids.append(bid)
            refs = sorted(box_refs.pop(bid, []))
            if bid not in box_sources:
                problems.append(f"canonical box {bid} not matched in PPTX")
            if len(refs) != len(set(refs)):
                problems.append(f"canonical box {bid} has duplicate refs: {refs}")
            boxes.append({"id": bid, "abbrev": abbrev, "label": label,
                          "refs": refs})
        out = {"id": comp["id"], "letter": comp["letter"],
               "title": comp["title"], "kind": comp["kind"],
               "panel": comp["panel"]}
        if "phaseLabel" in comp:
            out["phaseLabel"] = comp["phaseLabel"]
        out["boxes"] = boxes
        components_out.append(out)

    bands_out = []
    for band in BANDS:
        groups = []
        for grp in band["groups"]:
            boxes = []
            for suffix, abbrev in grp["boxes"]:
                bid = f"{grp['id']}.{suffix}"
                all_box_ids.append(bid)
                refs = sorted(box_refs.pop(bid, []))
                if bid not in box_sources:
                    problems.append(f"canonical box {bid} not matched in PPTX")
                if len(refs) != len(set(refs)):
                    problems.append(f"canonical box {bid} has duplicate refs: {refs}")
                boxes.append({"id": bid, "abbrev": abbrev,
                              "label": MOD_LABEL[abbrev], "refs": refs})
            groups.append({"id": grp["id"], "title": grp["title"],
                           "boxes": boxes})
        bands_out.append({"id": band["id"], "panel": band["panel"],
                          "groups": groups})

    if box_refs:  # refs resolved to a box id not in the inventory
        problems.append(f"extra ref-bearing boxes outside inventory: "
                        f"{sorted(box_refs)}")

    framework = {
        "meta": {"source": "draw main figure ver2.pptx", "refTotal": 189},
        "components": components_out,
        "bands": bands_out,
        "edges": [],
    }

    # ---- audit draft dump ----------------------------------------------
    draft = {
        "geometry": {**{k: geo[k] for k in "abcdefgh"},
                     "p2band": geo["p2band"], "p3band": geo["p3band"]},
        "labels": [{k: lb[k] for k in
                    ("text", "x", "y", "w", "h", "bucket", "canonical")}
                   | {"refs": sorted(lb.get("_refs", []))}
                   for lb in labels],
        "refboxes": [{"text": rb["text"], "refs": rb["refs"],
                      "x": rb["x"], "y": rb["y"],
                      "assigned_label": rb.get("label"),
                      "canonical": (rb.get("_target") or {}).get("canonical")}
                     for rb in refboxes],
        "ignored_labels": [{"text": lb["text"][:120], "x": lb["x"],
                            "y": lb["y"], "bucket": lb["bucket"]}
                           for lb in ignored],
        "totals": {
            "ref_occurrences": sum(len(rb["refs"]) for rb in refboxes),
            "labels": len(labels),
            "refboxes": len(refboxes),
        },
    }
    DRAFT_JSON.write_text(json.dumps(draft, indent=2, ensure_ascii=False))

    if problems:
        print("EXTRACTION FAILED — structural problems:", file=sys.stderr)
        for p in problems:
            print(f"  - {p}", file=sys.stderr)
        print(f"(draft dump written to {DRAFT_JSON} for inspection)",
              file=sys.stderr)
        sys.exit(1)

    OUT_JSON.parent.mkdir(parents=True, exist_ok=True)
    OUT_JSON.write_text(json.dumps(framework, indent=2, ensure_ascii=False))

    n_boxes = len(all_box_ids)
    n_refs = sum(len(rb["refs"]) for rb in refboxes)
    distinct = sorted({r for rb in refboxes for r in rb["refs"]})
    print(f"OK: wrote {OUT_JSON}")
    print(f"  canonical boxes: {n_boxes}; ref occurrences: {n_refs}; "
          f"distinct refs: {len(distinct)}")
    print(f"  ignored decorative/annotation labels: {len(ignored)} "
          f"(see {DRAFT_JSON.name})")


if __name__ == "__main__":
    main()
